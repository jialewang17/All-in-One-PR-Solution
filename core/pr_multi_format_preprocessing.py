#!/usr/bin/env python3
"""
公关传播内容多格式预处理脚本
支持PDF、Excel、CSV、Word、PPT、HTML、JSON、TXT等格式
"""

import os
import re
import json
import pandas as pd
from pathlib import Path
from docx import Document
from bs4 import BeautifulSoup
import PyPDF2
import pptx
import warnings
warnings.filterwarnings("ignore")

def read_pdf_file(file_path):
    """读取PDF文件"""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text_content = []
            
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text_content.append(page.extract_text())
            
            return '\n'.join(text_content)
    except Exception as e:
        print(f"Error reading PDF file {file_path}: {e}")
        return None

def read_excel_file(file_path):
    """读取Excel文件"""
    try:
        # 读取所有工作表
        excel_file = pd.ExcelFile(file_path)
        text_content = []
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            text_content.append(f"Sheet: {sheet_name}")
            text_content.append(df.to_string())
        
        return '\n'.join(text_content)
    except Exception as e:
        print(f"Error reading Excel file {file_path}: {e}")
        return None

def read_csv_file(file_path):
    """读取CSV文件"""
    try:
        df = pd.read_csv(file_path, encoding='utf-8')
        return df.to_string()
    except Exception as e:
        try:
            # 尝试其他编码
            df = pd.read_csv(file_path, encoding='gbk')
            return df.to_string()
        except Exception as e2:
            print(f"Error reading CSV file {file_path}: {e2}")
            return None

def read_docx_file(file_path):
    """读取Word文档"""
    try:
        doc = Document(file_path)
        text_content = []
        
        # 提取段落文本
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text.strip())
        
        # 提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content.append(" | ".join(row_text))
        
        return '\n'.join(text_content)
    except Exception as e:
        print(f"Error reading DOCX file {file_path}: {e}")
        return None

def read_pptx_file(file_path):
    """读取PowerPoint文件"""
    try:
        prs = pptx.Presentation(file_path)
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"Slide {slide_num}:")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
        
        return '\n'.join(text_content)
    except Exception as e:
        print(f"Error reading PPTX file {file_path}: {e}")
        return None

def read_html_file(file_path):
    """读取HTML文件"""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            html_content = file.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # 移除脚本和样式标签
        for script in soup(["script", "style"]):
            script.decompose()
        
        # 提取文本
        text_content = []
        for element in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'div']):
            if element.get_text().strip():
                text_content.append(element.get_text().strip())
        
        return '\n'.join(text_content)
    except Exception as e:
        print(f"Error reading HTML file {file_path}: {e}")
        return None

def read_json_file(file_path):
    """读取JSON文件"""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
        
        # 将JSON转换为可读文本
        if isinstance(data, dict):
            text_content = []
            for key, value in data.items():
                text_content.append(f"{key}: {value}")
            return '\n'.join(text_content)
        elif isinstance(data, list):
            return '\n'.join([str(item) for item in data])
        else:
            return str(data)
    except Exception as e:
        print(f"Error reading JSON file {file_path}: {e}")
        return None

def read_txt_file(file_path):
    """读取TXT文件"""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        try:
            # 尝试其他编码
            with open(file_path, 'r', encoding='gbk') as file:
                return file.read()
        except Exception as e2:
            print(f"Error reading TXT file {file_path}: {e2}")
            return None

def extract_text_from_content(content, file_type):
    """从内容中提取和格式化文本"""
    if not content:
        return None
    
    lines = content.split('\n')
    text_content = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # 根据文件类型和内容特征判断是否为标题
        is_title = False
        
        if file_type in ['pdf', 'docx', 'pptx']:
            # 对于结构化文档，判断标题
            if len(line) < 100 and (line.endswith('：') or line.endswith(':') or 
                                   '品牌' in line or '案例' in line or '策略' in line or 
                                   '传播' in line or '营销' in line or '分析' in line):
                is_title = True
        elif file_type == 'html':
            # HTML文档的标题判断
            if len(line) < 100 and ('品牌' in line or '案例' in line or '策略' in line):
                is_title = True
        elif file_type in ['excel', 'csv']:
            # 表格数据的处理
            if 'Sheet:' in line or line.startswith('Unnamed:'):
                is_title = True
        
        if is_title:
            text_content.append(f"Section: {line}")
        else:
            text_content.append(f"Content: {line}")
    
    return '\n'.join(text_content)

def save_text_to_file(text, output_path):
    """保存提取的文本到文件"""
    try:
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as file:
            file.write(text)
        print(f"Text saved to: {output_path}")
        return True
    except Exception as e:
        print(f"Error saving file {output_path}: {e}")
        return False

def process_multi_format_documents(input_dir="data/raw", output_dir="data/cleaned"):
    """处理多种格式的公关传播文档"""
    input_path = Path(input_dir)
    output_path = Path(output_dir)
    
    if not input_path.exists():
        print(f"Input directory {input_dir} does not exist")
        return
    
    # 创建输出目录
    output_path.mkdir(parents=True, exist_ok=True)
    
    # 支持的文件格式
    supported_formats = {
        '.pdf': read_pdf_file,
        '.xlsx': read_excel_file,
        '.xls': read_excel_file,
        '.csv': read_csv_file,
        '.docx': read_docx_file,
        '.doc': read_docx_file,
        '.pptx': read_pptx_file,
        '.ppt': read_pptx_file,
        '.html': read_html_file,
        '.htm': read_html_file,
        '.json': read_json_file,
        '.txt': read_txt_file
    }
    
    # 统计处理的文件
    processed_files = 0
    
    for file_path in input_path.iterdir():
        if file_path.is_file():
            file_ext = file_path.suffix.lower()
            
            if file_ext in supported_formats:
                print(f"\nProcessing: {file_path.name} ({file_ext})")
                
                # 读取文件内容
                content = supported_formats[file_ext](file_path)
                
                if content:
                    # 提取文本
                    text_content = extract_text_from_content(content, file_ext[1:])
                    
                    if text_content:
                        # 生成输出文件名
                        output_filename = file_path.stem + ".txt"
                        output_file_path = output_path / output_filename
                        
                        # 保存文本
                        if save_text_to_file(text_content, output_file_path):
                            print(f"✅ Successfully processed {file_path.name}")
                            processed_files += 1
                        else:
                            print(f"❌ Failed to process {file_path.name}")
                    else:
                        print(f"❌ No text content extracted from {file_path.name}")
                else:
                    print(f"❌ Failed to read {file_path.name}")
            else:
                print(f"⚠️ Unsupported format: {file_path.name} ({file_ext})")
    
    print(f"\n📊 处理完成！成功处理了 {processed_files} 个文件")
    print(f"支持的文件格式: {', '.join(supported_formats.keys())}")

if __name__ == "__main__":
    print("🚀 公关传播多格式文档预处理开始")
    print("="*60)
    print("支持格式: PDF, Excel, CSV, Word, PowerPoint, HTML, JSON, TXT")
    print("="*60)
    
    # 处理文档
    process_multi_format_documents()
    
    print("\n✅ 多格式预处理完成！")
    print("处理后的文件保存在 data/cleaned/ 目录中")


